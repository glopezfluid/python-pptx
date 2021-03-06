# encoding: utf-8

"""
Test suite for pptx.slide module
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.parts.presentation import PresentationPart
from pptx.parts.slide import SlideLayoutPart, SlideMasterPart, SlidePart
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.shapes.shapetree import (
    LayoutPlaceholders, LayoutShapes, MasterPlaceholders, MasterShapes,
    SlidePlaceholders, SlideShapes
)
from pptx.slide import (
    _BaseSlide, Slide, SlideLayout, SlideLayouts, SlideMaster, SlideMasters,
    Slides
)

from .unitutil.cxml import element, xml
from .unitutil.mock import call, class_mock, instance_mock, property_mock


class Describe_BaseSlide(object):

    def it_knows_its_name(self, name_get_fixture):
        base_slide, expected_value = name_get_fixture
        assert base_slide.name == expected_value

    def it_can_change_its_name(self, name_set_fixture):
        base_slide, new_value, expected_xml = name_set_fixture
        base_slide.name = new_value
        assert base_slide._element.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:sld/p:cSld',              ''),
        ('p:sld/p:cSld{name=Foobar}', 'Foobar'),
    ])
    def name_get_fixture(self, request):
        sld_cxml, expected_name = request.param
        base_slide = _BaseSlide(element(sld_cxml), None)
        return base_slide, expected_name

    @pytest.fixture(params=[
        ('p:sld/p:cSld',           'foo', 'p:sld/p:cSld{name=foo}'),
        ('p:sld/p:cSld{name=foo}', 'bar', 'p:sld/p:cSld{name=bar}'),
        ('p:sld/p:cSld{name=bar}', '',    'p:sld/p:cSld'),
        ('p:sld/p:cSld{name=bar}', None,  'p:sld/p:cSld'),
        ('p:sld/p:cSld',           '',    'p:sld/p:cSld'),
        ('p:sld/p:cSld',           None,  'p:sld/p:cSld'),
    ])
    def name_set_fixture(self, request):
        xSld_cxml, new_value, expected_cxml = request.param
        base_slide = _BaseSlide(element(xSld_cxml), None)
        expected_xml = xml(expected_cxml)
        return base_slide, new_value, expected_xml


class DescribeSlide(object):

    def it_is_a_BaseSlide_subclass(self, subclass_fixture):
        slide = subclass_fixture
        assert isinstance(slide, _BaseSlide)

    def it_knows_its_slide_id(self, slide_id_fixture):
        slide, expected_value = slide_id_fixture
        assert slide.slide_id == expected_value

    def it_provides_access_to_its_shapes(self, shapes_fixture):
        slide, SlideShapes_, spTree, shapes_ = shapes_fixture
        shapes = slide.shapes
        SlideShapes_.assert_called_once_with(spTree, slide)
        assert shapes is shapes_

    def it_provides_access_to_its_placeholders(self, placeholders_fixture):
        slide, SlidePlaceholders_, spTree, placeholders_ = (
            placeholders_fixture
        )
        placeholders = slide.placeholders
        SlidePlaceholders_.assert_called_once_with(spTree, slide)
        assert placeholders is placeholders_

    def it_provides_access_to_its_slide_layout(self, layout_fixture):
        slide, slide_layout_ = layout_fixture
        assert slide.slide_layout is slide_layout_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def layout_fixture(self, slide_layout_, part_prop_, slide_part_):
        slide = Slide(None, None)
        part_prop_.return_value = slide_part_
        slide_part_.slide_layout = slide_layout_
        return slide, slide_layout_

    @pytest.fixture
    def placeholders_fixture(self, SlidePlaceholders_, placeholders_):
        sld = element('p:sld/p:cSld/p:spTree')
        slide = Slide(sld, None)
        spTree = sld.xpath('//p:spTree')[0]
        return slide, SlidePlaceholders_, spTree, placeholders_

    @pytest.fixture
    def shapes_fixture(self, SlideShapes_, shapes_):
        sld = element('p:sld/p:cSld/p:spTree')
        spTree = sld.xpath('//p:spTree')[0]
        slide = Slide(sld, None)
        return slide, SlideShapes_, spTree, shapes_

    @pytest.fixture
    def slide_id_fixture(self, part_prop_, slide_part_):
        slide = Slide(None, None)
        slide_id = 256
        slide_part_.slide_id = slide_id
        return slide, slide_id

    @pytest.fixture
    def subclass_fixture(self):
        return Slide(None, None)

    # fixture components -----------------------------------

    @pytest.fixture
    def part_prop_(self, request, slide_part_):
        return property_mock(
            request, Slide, 'part', return_value=slide_part_
        )

    @pytest.fixture
    def placeholders_(self, request):
        return instance_mock(request, SlidePlaceholders)

    @pytest.fixture
    def SlidePlaceholders_(self, request, placeholders_):
        return class_mock(
            request, 'pptx.slide.SlidePlaceholders',
            return_value=placeholders_
        )

    @pytest.fixture
    def SlideShapes_(self, request, shapes_):
        return class_mock(
            request, 'pptx.slide.SlideShapes', return_value=shapes_
        )

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, SlideShapes)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_part_(self, request):
        return instance_mock(request, SlidePart)


class DescribeSlides(object):

    def it_supports_indexed_access(self, getitem_fixture):
        slides, prs_part_, rId, slide_ = getitem_fixture
        slide = slides[0]
        prs_part_.related_slide.assert_called_once_with(rId)
        assert slide is slide_

    def it_raises_on_slide_index_out_of_range(self, getitem_raises_fixture):
        slides = getitem_raises_fixture
        with pytest.raises(IndexError):
            slides[2]

    def it_knows_the_index_of_a_slide_it_contains(self, index_fixture):
        slides, slide, expected_value = index_fixture
        index = slides.index(slide)
        assert index == expected_value

    def it_raises_on_slide_not_in_collection(self, raises_fixture):
        slides, slide = raises_fixture
        with pytest.raises(ValueError):
            slides.index(slide)

    def it_can_iterate_its_slides(self, iter_fixture):
        slides, related_slide_, calls, expected_value = iter_fixture
        slide_lst = [s for s in slides]
        assert related_slide_.call_args_list == calls
        assert slide_lst == expected_value

    def it_supports_len(self, len_fixture):
        slides, expected_value = len_fixture
        assert len(slides) == expected_value

    def it_can_add_a_new_slide(self, add_fixture):
        slides, slide_layout_, part_ = add_fixture[:3]
        clone_layout_placeholders_, expected_xml, slide_ = add_fixture[3:]

        slide = slides.add_slide(slide_layout_)

        part_.add_slide.assert_called_once_with(slide_layout_)
        clone_layout_placeholders_.assert_called_once_with(slide_layout_)
        assert slides._sldIdLst.xml == expected_xml
        assert slide is slide_

    def it_finds_a_slide_by_slide_id(self, get_fixture):
        slides, slide_id, default, prs_part_, expected_value = get_fixture
        slide = slides.get(slide_id, default)
        prs_part_.get_slide.assert_called_once_with(slide_id)
        assert slide is expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def add_fixture(self, slide_layout_, part_prop_, slide_):
        slides = Slides(element('p:sldIdLst/p:sldId{r:id=rId1}'), None)
        part_ = part_prop_.return_value
        clone_layout_placeholders_ = slide_.shapes.clone_layout_placeholders
        expected_xml = xml(
            'p:sldIdLst/(p:sldId{r:id=rId1},p:sldId{r:id=rId2,id=256})'
        )
        part_.add_slide.return_value = 'rId2', slide_
        return (
            slides, slide_layout_, part_, clone_layout_placeholders_,
            expected_xml, slide_
        )

    @pytest.fixture(params=[True, False])
    def get_fixture(self, request, part_prop_, prs_part_, slide_):
        found = request.param
        slides = Slides(None, None)
        slide_id, default = 256, 'foobar'
        expected_value = slide_ if found else default
        prs_part_.get_slide.return_value = slide_ if found else None
        return slides, slide_id, default, prs_part_, expected_value

    @pytest.fixture
    def getitem_fixture(self, prs_part_, slide_, part_prop_):
        sldIdLst = element('p:sldIdLst/p:sldId{r:id=rId1}')
        slides = Slides(sldIdLst, None)
        prs_part_.related_slide.return_value = slide_
        return slides, prs_part_, 'rId1', slide_

    @pytest.fixture
    def getitem_raises_fixture(self):
        sldIdLst = element('p:sldIdLst/p:sldId{r:id=rId1}')
        slides = Slides(sldIdLst, None)
        return slides

    @pytest.fixture(params=[0, 1])
    def index_fixture(self, request, part_prop_):
        idx = request.param
        sldIdLst = element('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})')
        slides = Slides(sldIdLst, None)
        _slides = [
            Slide(element('p:sld'), None),
            Slide(element('p:sld'), None)
        ]
        part_prop_.return_value.related_slide.side_effect = (
            _slides
        )
        return slides, _slides[idx], idx

    @pytest.fixture
    def iter_fixture(self, part_prop_, slide_):
        sldIdLst = element('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})')
        slides = Slides(sldIdLst, None)
        related_slide_ = part_prop_.return_value.related_slide
        related_slide_.return_value = slide_
        calls = [call('a'), call('b')]
        _slides = [slide_, slide_]
        return slides, related_slide_, calls, _slides

    @pytest.fixture(params=[
        ('p:sldIdLst',                                   0),
        ('p:sldIdLst/p:sldId{r:id=a}',                   1),
        ('p:sldIdLst/(p:sldId{r:id=a},p:sldId{r:id=b})', 2),
    ])
    def len_fixture(self, request):
        sldIdLst_cxml, expected_value = request.param
        slides = Slides(element(sldIdLst_cxml), None)
        return slides, expected_value

    @pytest.fixture
    def raises_fixture(self):
        slides = Slides(element('p:sldIdLst'), None)
        slide = Slide(element('p:sld'), None)
        return slides, slide

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_prop_(self, request, prs_part_):
        return property_mock(request, Slides, 'part', return_value=prs_part_)

    @pytest.fixture
    def prs_part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)


class DescribeSlideLayout(object):

    def it_is_a_BaseSlide_subclass(self, subclass_fixture):
        slide = subclass_fixture
        assert isinstance(slide, _BaseSlide)

    def it_provides_access_to_its_slide_master(self, master_fixture):
        slide_layout, slide_master_ = master_fixture
        assert slide_layout.slide_master is slide_master_

    def it_provides_access_to_its_placeholders(self, placeholders_fixture):
        slide_layout, LayoutPlaceholders_, spTree, placeholders_ = (
            placeholders_fixture
        )
        placeholders = slide_layout.placeholders
        LayoutPlaceholders_.assert_called_once_with(spTree, slide_layout)
        assert placeholders is placeholders_

    def it_provides_access_to_its_shapes(self, shapes_fixture):
        slide_layout, LayoutShapes_, spTree, shapes_ = shapes_fixture
        shapes = slide_layout.shapes
        LayoutShapes_.assert_called_once_with(spTree, slide_layout)
        assert shapes is shapes_

    def it_can_iterate_its_clonable_placeholders(self, cloneable_fixture):
        slide_layout, expected_placeholders = cloneable_fixture
        cloneable = list(slide_layout.iter_cloneable_placeholders())
        assert cloneable == expected_placeholders

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ((PP_PLACEHOLDER.TITLE,        PP_PLACEHOLDER.BODY),   (0, 1)),
        ((PP_PLACEHOLDER.TITLE,        PP_PLACEHOLDER.DATE),   (0,)),
        ((PP_PLACEHOLDER.FOOTER,       PP_PLACEHOLDER.OBJECT), (1,)),
        ((PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.FOOTER), ()),
    ])
    def cloneable_fixture(
            self, request, placeholders_prop_, placeholder_, placeholder_2_):
        ph_types, expected_indices = request.param
        slide_layout = SlideLayout(None, None)
        placeholder_.ph_type, placeholder_2_.ph_type = ph_types
        _placeholders = (placeholder_, placeholder_2_)
        expected_placeholders = [
            _placeholders[idx] for idx in expected_indices
        ]
        placeholders_prop_.return_value = _placeholders
        return slide_layout, expected_placeholders

    @pytest.fixture
    def master_fixture(self, slide_master_, part_prop_):
        slide_layout = SlideLayout(None, None)
        part_prop_.return_value.slide_master = slide_master_
        return slide_layout, slide_master_

    @pytest.fixture
    def placeholders_fixture(self, LayoutPlaceholders_, placeholders_):
        sldLayout = element('p:sldLayout/p:cSld/p:spTree')
        slide_layout = SlideLayout(sldLayout, None)
        spTree = sldLayout.xpath('//p:spTree')[0]
        return slide_layout, LayoutPlaceholders_, spTree, placeholders_

    @pytest.fixture
    def shapes_fixture(self, LayoutShapes_, shapes_):
        sldLayout = element('p:sldLayout/p:cSld/p:spTree')
        slide_layout = SlideLayout(sldLayout, None)
        spTree = sldLayout.xpath('//p:spTree')[0]
        return slide_layout, LayoutShapes_, spTree, shapes_

    @pytest.fixture
    def subclass_fixture(self):
        return SlideLayout(None, None)

    # fixture components -----------------------------------

    @pytest.fixture
    def LayoutPlaceholders_(self, request, placeholders_):
        return class_mock(
            request, 'pptx.slide.LayoutPlaceholders',
            return_value=placeholders_
        )

    @pytest.fixture
    def LayoutShapes_(self, request, shapes_):
        return class_mock(
            request, 'pptx.slide.LayoutShapes', return_value=shapes_
        )

    @pytest.fixture
    def part_prop_(self, request, slide_layout_part_):
        return property_mock(
            request, SlideLayout, 'part', return_value=slide_layout_part_
        )

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def placeholder_2_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def placeholders_(self, request):
        return instance_mock(request, LayoutPlaceholders)

    @pytest.fixture
    def placeholders_prop_(self, request, placeholders_):
        return property_mock(
            request, SlideLayout, 'placeholders', return_value=placeholders_
        )

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, LayoutShapes)

    @pytest.fixture
    def slide_layout_part_(self, request):
        return instance_mock(request, SlideLayoutPart)

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)


class DescribeSlideLayouts(object):

    def it_supports_len(self, len_fixture):
        slide_layouts, expected_value = len_fixture
        assert len(slide_layouts) == expected_value

    def it_can_iterate_its_slide_layouts(self, iter_fixture):
        slide_layouts, related_slide_layout_ = iter_fixture[:2]
        calls, expected_value = iter_fixture[2:]
        slide_layout_lst = [sl for sl in slide_layouts]
        assert related_slide_layout_.call_args_list == calls
        assert slide_layout_lst == expected_value

    def it_supports_indexed_access(self, getitem_fixture):
        slide_layouts, part_, slide_layout_, rId = getitem_fixture
        slide_layout = slide_layouts[0]
        part_.related_slide_layout.assert_called_once_with(rId)
        assert slide_layout is slide_layout_

    def it_raises_on_index_out_of_range(self, getitem_raises_fixture):
        slides = getitem_raises_fixture
        with pytest.raises(IndexError):
            slides[1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def getitem_fixture(self, part_, slide_layout_, part_prop_):
        slide_layouts = SlideLayouts(
            element('p:sldLayoutIdLst/p:sldLayoutId{r:id=rId1}'), None
        )
        part_.related_slide_layout.return_value = slide_layout_
        return slide_layouts, part_, slide_layout_, 'rId1'

    @pytest.fixture
    def getitem_raises_fixture(self, part_prop_):
        return SlideLayouts(
            element('p:sldLayoutIdLst/p:sldLayoutId{r:id=rId1}'), None
        )

    @pytest.fixture
    def iter_fixture(self, part_prop_, part_):
        sldLayoutIdLst = element(
            'p:sldLayoutIdLst/(p:sldLayoutId{r:id=a},p:sldLayoutId{r:id=b})'
        )
        slide_layouts = SlideLayouts(sldLayoutIdLst, None)
        _slide_layouts = [
            SlideLayout(element('p:sldLayout'), None),
            SlideLayout(element('p:sldLayout'), None),
        ]
        related_slide_layout_ = part_.related_slide_layout
        related_slide_layout_.side_effect = _slide_layouts
        calls = [call('a'), call('b')]
        return slide_layouts, related_slide_layout_, calls, _slide_layouts

    @pytest.fixture(params=[
        ('p:sldLayoutIdLst',                               0),
        ('p:sldLayoutIdLst/p:sldLayoutId',                 1),
        ('p:sldLayoutIdLst/(p:sldLayoutId,p:sldLayoutId)', 2),
    ])
    def len_fixture(self, request):
        sldLayoutIdLst_cxml, expected_value = request.param
        slide_layouts = SlideLayouts(element(sldLayoutIdLst_cxml), None)
        return slide_layouts, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_(self, request):
        return instance_mock(request, SlideMasterPart)

    @pytest.fixture
    def part_prop_(self, request, part_):
        return property_mock(
            request, SlideLayouts, 'part', return_value=part_
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)


class DescribeSlideMaster(object):

    def it_is_a_BaseSlide_subclass(self, subclass_fixture):
        slide = subclass_fixture
        assert isinstance(slide, _BaseSlide)

    def it_provides_access_to_its_placeholders(self, placeholders_fixture):
        slide_master, MasterPlaceholders_, spTree, placeholders_ = (
            placeholders_fixture
        )
        placeholders = slide_master.placeholders
        MasterPlaceholders_.assert_called_once_with(spTree, slide_master)
        assert placeholders is placeholders_

    def it_provides_access_to_its_shapes(self, shapes_fixture):
        slide_master, MasterShapes_, spTree, shapes_ = shapes_fixture
        shapes = slide_master.shapes
        MasterShapes_.assert_called_once_with(spTree, slide_master)
        assert shapes is shapes_

    def it_provides_access_to_its_slide_layouts(self, layouts_fixture):
        slide_master, SlideLayouts_, sldLayoutIdLst, slide_layouts_ = (
            layouts_fixture
        )
        slide_layouts = slide_master.slide_layouts
        SlideLayouts_.assert_called_once_with(sldLayoutIdLst, slide_master)
        assert slide_layouts is slide_layouts_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def layouts_fixture(self, SlideLayouts_, slide_layouts_):
        sldMaster = element('p:sldMaster/p:sldLayoutIdLst')
        slide_master = SlideMaster(sldMaster, None)
        sldMasterIdLst = sldMaster.sldLayoutIdLst
        return slide_master, SlideLayouts_, sldMasterIdLst, slide_layouts_

    @pytest.fixture
    def placeholders_fixture(self, MasterPlaceholders_, placeholders_):
        sldMaster = element('p:sldMaster/p:cSld/p:spTree')
        slide_master = SlideMaster(sldMaster, None)
        spTree = sldMaster.xpath('//p:spTree')[0]
        return slide_master, MasterPlaceholders_, spTree, placeholders_

    @pytest.fixture
    def shapes_fixture(self, MasterShapes_, shapes_):
        sldMaster = element('p:sldMaster/p:cSld/p:spTree')
        slide_master = SlideMaster(sldMaster, None)
        spTree = sldMaster.xpath('//p:spTree')[0]
        return slide_master, MasterShapes_, spTree, shapes_

    @pytest.fixture
    def subclass_fixture(self):
        return SlideMaster(None, None)

    # fixture components -----------------------------------

    @pytest.fixture
    def MasterPlaceholders_(self, request, placeholders_):
        return class_mock(
            request, 'pptx.slide.MasterPlaceholders',
            return_value=placeholders_
        )

    @pytest.fixture
    def MasterShapes_(self, request, shapes_):
        return class_mock(
            request, 'pptx.slide.MasterShapes', return_value=shapes_
        )

    @pytest.fixture
    def placeholders_(self, request):
        return instance_mock(request, MasterPlaceholders)

    @pytest.fixture
    def shapes_(self, request):
        return instance_mock(request, MasterShapes)

    @pytest.fixture
    def SlideLayouts_(self, request, slide_layouts_):
        return class_mock(
            request, 'pptx.slide.SlideLayouts', return_value=slide_layouts_
        )

    @pytest.fixture
    def slide_layouts_(self, request):
        return instance_mock(request, SlideLayouts)


class DescribeSlideMasters(object):

    def it_knows_how_many_masters_it_contains(self, len_fixture):
        slide_masters, expected_value = len_fixture
        assert len(slide_masters) == expected_value

    def it_can_iterate_the_slide_masters(self, iter_fixture):
        slide_masters, related_slide_master_, calls, expected_values = (
            iter_fixture
        )
        _slide_masters = [sm for sm in slide_masters]
        assert related_slide_master_.call_args_list == calls
        assert _slide_masters == expected_values

    def it_supports_indexed_access(self, getitem_fixture):
        slide_masters, part_, slide_master_, rId = getitem_fixture
        slide_master = slide_masters[0]
        part_.related_slide_master.assert_called_once_with(rId)
        assert slide_master is slide_master_

    def it_raises_on_index_out_of_range(self, getitem_raises_fixture):
        slides = getitem_raises_fixture
        with pytest.raises(IndexError):
            slides[1]

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def getitem_fixture(self, part_, slide_master_, part_prop_):
        slide_masters = SlideMasters(
            element('p:sldMasterIdLst/p:sldMasterId{r:id=rId1}'), None
        )
        part_.related_slide_master.return_value = slide_master_
        return slide_masters, part_, slide_master_, 'rId1'

    @pytest.fixture
    def getitem_raises_fixture(self, part_prop_):
        return SlideMasters(
            element('p:sldMasterIdLst/p:sldMasterId{r:id=rId1}'), None
        )

    @pytest.fixture
    def iter_fixture(self, part_prop_):
        sldMasterIdLst = element(
            'p:sldMasterIdLst/(p:sldMasterId{r:id=a},p:sldMasterId{r:id=b})'
        )
        slide_masters = SlideMasters(sldMasterIdLst, None)
        related_slide_master_ = part_prop_.return_value.related_slide_master
        calls = [call('a'), call('b')]
        _slide_masters = [
            SlideMaster(element('p:sldMaster'), None),
            SlideMaster(element('p:sldMaster'), None)
        ]
        related_slide_master_.side_effect = _slide_masters
        return slide_masters, related_slide_master_, calls, _slide_masters

    @pytest.fixture(params=[
        ('p:sldMasterIdLst',                               0),
        ('p:sldMasterIdLst/p:sldMasterId',                 1),
        ('p:sldMasterIdLst/(p:sldMasterId,p:sldMasterId)', 2),
    ])
    def len_fixture(self, request):
        sldMasterIdLst_cxml, expected_value = request.param
        slide_masters = SlideMasters(element(sldMasterIdLst_cxml), None)
        return slide_masters, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def part_prop_(self, request, part_):
        return property_mock(
            request, SlideMasters, 'part', return_value=part_
        )

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)
